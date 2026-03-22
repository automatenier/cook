import sys
import json
from pptx import Presentation
from pptx.util import Inches

def create_pptx(data, output_path):
    prs = Presentation()
    
    # Title Slide
    if "title" in data:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = data["title"]
        if "subtitle" in data:
            subtitle.text = data["subtitle"]

    # Content Slides
    for slide_data in data.get("slides", []):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data.get("title", "")
        
        body = slide.placeholders[1]
        tf = body.text_frame
        
        for i, point in enumerate(slide_data.get("content", [])):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: py -3 create_pptx.py <input_json_path> <output_pptx_path>")
        sys.exit(1)
        
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)
        
    create_pptx(data, output_path)
