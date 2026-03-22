import sys
import json
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def draw_mindmap(slide, data):
    """Draws a simple radial mindmap."""
    width = Inches(10)
    height = Inches(7.5)
    center_x = width / 2
    center_y = height / 2
    
    # Draw Central Topic
    root_text = data.get("root", "Central Topic")
    root_w, root_h = Inches(2), Inches(1)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        center_x - root_w/2, center_y - root_h/2, root_w, root_h
    )
    shape.text = root_text
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 112, 192) # Blue
    
    # Draw Nodes
    nodes = data.get("nodes", [])
    num_nodes = len(nodes)
    if num_nodes == 0: return

    radius = Inches(2.5)
    for i, node in enumerate(nodes):
        angle = (2 * math.pi * i) / num_nodes
        x = center_x + radius * math.cos(angle)
        y = center_y + radius * math.sin(angle)
        
        node_w, node_h = Inches(1.5), Inches(0.8)
        node_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            x - node_w/2, y - node_h/2, node_w, node_h
        )
        node_shape.text = node
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = RGBColor(146, 208, 80) # Green
        
        # Draw connector
        line = slide.shapes.add_connector(
            1, center_x, center_y, x, y
        )

def create_pptx(data, output_path):
    prs = Presentation()
    
    # Title Slide
    if "title" in data:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = data["title"]
        if "subtitle" in data:
            slide.placeholders[1].text = data["subtitle"]

    # Content & Mindmap Slides
    for slide_data in data.get("slides", []):
        stype = slide_data.get("type", "content")
        
        if stype == "mindmap":
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_shape.text_frame.text = slide_data.get("title", "Mind Map")
            draw_mindmap(slide, slide_data.get("map_data", {}))
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", "")
            tf = slide.placeholders[1].text_frame
            for i, point in enumerate(slide_data.get("content", [])):
                if i == 0: tf.text = point
                else: tf.add_paragraph().text = point
                
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: py -3 create_pptx.py <input_json_path> <output_pptx_path>")
        sys.exit(1)
    input_path, output_path = sys.argv[1], sys.argv[2]
    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    create_pptx(data, output_path)
