import sys
from pptx import Presentation

def read_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: py -3 read_pptx.py <pptx_path>")
        sys.exit(1)
    
    print(read_pptx(sys.argv[1]))
